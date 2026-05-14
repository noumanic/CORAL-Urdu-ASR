"""
CORAL Final Year Project - Animated PPTX Slide Builder
=======================================================
Builds a 16:9, modern Canva-style animated PPTX summarising the entire
CORAL Urdu ASR R&D project (merged FYP1 + FYP2).

Run:  python build_slides.py
Out:  ../report/CORAL-FYP-Presentation.pptx
"""

from __future__ import annotations
import os
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree

# ---------------------------------------------------------------------------
# Theme - modern dark-on-light coral/ocean palette
# ---------------------------------------------------------------------------
THEME = {
    "bg":        RGBColor(0x0B, 0x1E, 0x33),   # deep navy
    "bg_light":  RGBColor(0xF5, 0xF7, 0xFA),   # off-white
    "primary":   RGBColor(0xFF, 0x6B, 0x6B),   # coral red
    "secondary": RGBColor(0x4E, 0xCD, 0xC4),   # teal
    "accent":    RGBColor(0xFF, 0xD9, 0x3D),   # warm yellow
    "violet":    RGBColor(0x8B, 0x5C, 0xF6),   # violet
    "green":     RGBColor(0x10, 0xB9, 0x81),   # emerald
    "orange":    RGBColor(0xF5, 0x9E, 0x0B),   # amber
    "blue":      RGBColor(0x3B, 0x82, 0xF6),   # cobalt
    "text":      RGBColor(0xE5, 0xE7, 0xEB),   # near white
    "text_dim":  RGBColor(0x9C, 0xA3, 0xAF),   # slate
    "text_dark": RGBColor(0x1F, 0x29, 0x37),   # dark slate for light slides
    "card":      RGBColor(0x16, 0x2A, 0x42),   # card on dark
    "card_lt":   RGBColor(0xFF, 0xFF, 0xFF),
    "line":      RGBColor(0x37, 0x4B, 0x66),
}

FIG_DIR  = Path(__file__).parent / "latex" / "ThesisFigs"
FRONT_DIR = Path(__file__).parent.parent.parent / "coral_app" / "frontend" / "my-next-app" / "public"
OUT      = Path(__file__).parent / "CORAL-FYP-Presentation.pptx"

# ---------------------------------------------------------------------------
# Presentation setup (16:9)
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_rect(slide, x, y, w, h, fill, line=None, shadow=False,
             corner=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    if not shadow:
        s.shadow.inherit = False
    return s


def add_text(slide, text, x, y, w, h, *, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    paras = text.split("\n")
    for i, line in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color
    return tb


def add_pill(slide, text, x, y, w, h, fill, fg, size=11, bold=True):
    p = add_rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p.adjustments[0] = 0.5
    tf = p.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg
    return p


def add_card(slide, x, y, w, h, fill=None, line=None, corner=0.08):
    fill = fill if fill is not None else THEME["card"]
    s = add_rect(slide, x, y, w, h, fill, line=line,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    s.adjustments[0] = corner
    return s


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        return None
    kwargs = {}
    if w is not None: kwargs["width"] = w
    if h is not None: kwargs["height"] = h
    return slide.shapes.add_picture(str(path), x, y, **kwargs)


def add_decoration(slide, dark=True):
    """Decorative coral/teal blobs in corners."""
    # Bottom-right coral blob
    s1 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(11.6), Inches(6.0), Inches(2.6), Inches(2.6))
    s1.line.fill.background()
    s1.fill.solid()
    s1.fill.fore_color.rgb = THEME["primary"]
    s1.shadow.inherit = False
    _set_transparency(s1, 75)

    # Top-left teal blob
    s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(-1.0), Inches(-1.0), Inches(2.4), Inches(2.4))
    s2.line.fill.background()
    s2.fill.solid()
    s2.fill.fore_color.rgb = THEME["secondary"]
    s2.shadow.inherit = False
    _set_transparency(s2, 80)


def _set_transparency(shape, pct):
    """Set fill transparency (0-100%)."""
    sp = shape.fill.fore_color._xFill
    # Get solidFill element
    solidFill = sp
    # Find srgbClr child
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # Remove existing alpha
    for a in srgb.findall(qn("a:alpha")):
        srgb.remove(a)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int((100 - pct) * 1000)))


def add_header(slide, kicker, title, *, dark=True):
    fg_kicker = THEME["primary"]
    fg_title  = THEME["text"] if dark else THEME["text_dark"]
    add_text(slide, kicker.upper(), Inches(0.6), Inches(0.45),
             Inches(12), Inches(0.4), size=13, bold=True,
             color=fg_kicker, font="Calibri")
    add_text(slide, title, Inches(0.6), Inches(0.85),
             Inches(12), Inches(0.95), size=34, bold=True,
             color=fg_title, font="Calibri")
    # underline accent
    line = slide.shapes.add_connector(1, Inches(0.6), Inches(1.85),
                                       Inches(1.8), Inches(1.85))
    line.line.color.rgb = THEME["primary"]
    line.line.width = Pt(4)


def add_footer(slide, page_num, total, *, dark=True):
    fg = THEME["text_dim"] if dark else THEME["text_dark"]
    add_text(slide, "CORAL  |  Consensus-Based Refinement and Output Realignment",
             Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             size=10, color=fg)
    add_text(slide, f"{page_num} / {total}",
             Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.3),
             size=10, color=fg, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Animation XML injection
# ---------------------------------------------------------------------------
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

def _shape_id(sp):
    return int(sp.shape_id)


def inject_fade_animations(slide, shape_ids, *, delay_ms=200, base_ms=300):
    """No-op (kept for call-site compatibility).

    Originally injected a per-shape sequential fade-in entrance, but each
    shape became its own click trigger in PowerPoint's slideshow,
    forcing many Enter-presses to advance a single slide.  We now rely
    solely on slide transitions for motion; shapes appear immediately
    when the slide is shown.  Content is unchanged.
    """
    if not shape_ids:
        return
    sld = slide._element
    for t in sld.findall(qn("p:timing")):
        sld.remove(t)


def _seq_for_shapes(shape_ids, delay_ms, base_ms):
    """Build the child timeline nodes for a sequence of fade-in entrances."""
    out = []
    cur_id = 3
    for idx, sid in enumerate(shape_ids):
        d = delay_ms if idx > 0 else 0
        out.append(f"""
        <p:par>
          <p:cTn id="{cur_id}" fill="hold">
            <p:stCondLst><p:cond delay="{d}"/></p:stCondLst>
            <p:childTnLst>
              <p:par>
                <p:cTn id="{cur_id+1}" fill="hold">
                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="{cur_id+2}" presetID="10" presetClass="entr"
                             presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                        <p:childTnLst>
                          <p:set>
                            <p:cBhvr>
                              <p:cTn id="{cur_id+3}" dur="1" fill="hold">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                              </p:cTn>
                              <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                              <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                            </p:cBhvr>
                            <p:to><p:strVal val="visible"/></p:to>
                          </p:set>
                          <p:anim calcmode="lin" valueType="num">
                            <p:cBhvr additive="base">
                              <p:cTn id="{cur_id+4}" dur="{base_ms}" fill="hold"/>
                              <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                              <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                            </p:cBhvr>
                            <p:tavLst>
                              <p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
                              <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
                            </p:tavLst>
                          </p:anim>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
              </p:par>
            </p:childTnLst>
          </p:cTn>
        </p:par>
        """)
        cur_id += 6
    return "".join(out)


def set_slide_transition(slide, kind="fade"):
    """Add slide transition (fade/push)."""
    sld = slide._element
    for t in sld.findall(qn("p:transition")):
        sld.remove(t)
    transition_xml = f"""
<p:transition xmlns:p="{P_NS}" spd="med">
  <p:{kind}/>
</p:transition>
"""
    sld.append(etree.fromstring(transition_xml))


def collect_anim(slide, exclude=None):
    """Return shape_ids of all top-level shapes except backgrounds/decorations."""
    exclude = exclude or set()
    ids = []
    for sp in slide.shapes:
        try:
            sid = int(sp.shape_id)
        except Exception:
            continue
        if sid in exclude:
            continue
        ids.append(sid)
    return ids


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
TOTAL_SLIDES = 32   # updated below

def page(num, total=TOTAL_SLIDES, dark=True):
    s = prs.slides.add_slide(BLANK)
    bg = add_bg(s, THEME["bg"] if dark else THEME["bg_light"])
    return s, bg


# ---------------------------------------------------------------------------
# Programmatic diagram builders (drawn from python-pptx shapes)
# ---------------------------------------------------------------------------
def _chip(slide, text, x, y, w, h, fill, fg=None, size=10, bold=True,
          corner=0.18):
    """Rounded-rectangle word chip with centered label."""
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.adjustments[0] = corner
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = fill
    c.line.width = Pt(0.5)
    c.shadow.inherit = False
    add_text(slide, text, x, y, w, h,
             size=size, bold=bold,
             color=fg if fg else RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def _line_between(slide, x1, y1, x2, y2, color, width=1.0, dashed=False):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        ln = line.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    return line


def draw_splitmerge_diagram(slide, x, y, w, h):
    """Hand-drawn Stage 1 split-merge alignment diagram."""
    # background card
    bg = add_card(slide, x, y, w, h,
                  fill=RGBColor(0x12, 0x1F, 0x35), line=THEME["line"],
                  corner=0.04)

    pad = Inches(0.25)
    inner_w = w - 2 * pad
    row_h = Inches(0.5)
    row_gap = Inches(0.55)

    # Row labels
    label_w = Inches(1.4)
    chip_area_x = x + pad + label_w + Inches(0.1)
    chip_area_w = inner_w - label_w - Inches(0.1)

    # Row 1 — Source (whisper-large): 4 chips, includes a MERGED chip
    # Row 2 — Companion 1 (seamless): 5 chips
    # Row 3 — Companion 2 (wav2vec):  4 chips, includes another merge

    rows = [
        ("whisper-L · src", THEME["accent"],
         [("A", THEME["green"]),
          ("B-C", THEME["orange"]),    # merged token
          ("D", THEME["green"]),
          ("E", THEME["green"])]),
        ("seamless-L",      THEME["blue"],
         [("A",  THEME["green"]),
          ("B",  THEME["primary"]),    # split, half 1
          ("C",  THEME["primary"]),    # split, half 2
          ("D",  THEME["green"]),
          ("E",  THEME["green"])]),
        ("wav2vec2",        THEME["violet"],
         [("A",   THEME["green"]),
          ("BC",  THEME["orange"]),    # merged like source
          ("D",   THEME["green"]),
          ("E",   THEME["green"])]),
    ]

    # title strip
    add_text(slide, "Alignment input — three model attempts",
             x + pad, y + Inches(0.10), inner_w, Inches(0.30),
             size=11, bold=True, color=THEME["text"])

    yy = y + Inches(0.55)
    for label, label_color, chips in rows:
        # label box
        lbl = add_rect(slide, x + pad, yy, label_w, row_h,
                       THEME["card"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        lbl.adjustments[0] = 0.20
        lbl.line.color.rgb = label_color
        lbl.line.width = Pt(1.2)
        add_text(slide, label, x + pad, yy, label_w, row_h,
                 size=10, bold=True, color=label_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # chips
        n = len(chips)
        gap = Inches(0.08)
        chip_w = (chip_area_w - gap * (n - 1)) / n
        cx = chip_area_x
        for txt, col in chips:
            _chip(slide, txt, cx, yy, chip_w, row_h, fill=col, size=12)
            cx += chip_w + gap
        yy += row_h + Inches(0.22)

    # Bottom alignment classification chips
    yy_class = yy + Inches(0.05)
    add_text(slide, "Per-chunk classification",
             x + pad, yy_class, inner_w, Inches(0.3),
             size=10, bold=True, color=THEME["secondary"])

    classes = [
        ("SAME",  "A=A=A",       THEME["green"]),
        ("SPLIT", "B-C ↔ B,C",   THEME["primary"]),
        ("MERGE", "B-C ↔ BC",    THEME["orange"]),
        ("SAME",  "D=D=D",       THEME["green"]),
        ("SAME",  "E=E=E",       THEME["green"]),
    ]
    n = len(classes)
    gap = Inches(0.10)
    cw = (inner_w - gap * (n - 1)) / n
    cy = yy_class + Inches(0.32)
    cx = x + pad
    for tag, desc, col in classes:
        # outer chip
        box = add_card(slide, cx, cy, cw, Inches(0.78),
                       fill=THEME["card"], line=col, corner=0.10)
        box.line.width = Pt(1.5)
        add_pill(slide, tag, cx + Inches(0.10), cy + Inches(0.08),
                 cw - Inches(0.2), Inches(0.30),
                 col, RGBColor(0xFF, 0xFF, 0xFF), size=10)
        add_text(slide, desc, cx, cy + Inches(0.40),
                 cw, Inches(0.32),
                 size=10, color=THEME["text"], align=PP_ALIGN.CENTER,
                 font="Consolas")
        cx += cw + gap


def draw_oov_bktree_diagram(slide, x, y, w, h):
    """Hand-drawn Stage 2 OOV detection + BK-tree diagram."""
    bg = add_card(slide, x, y, w, h,
                  fill=RGBColor(0x12, 0x1F, 0x35), line=THEME["line"],
                  corner=0.04)

    pad = Inches(0.25)
    inner_x = x + pad
    inner_w = w - 2 * pad

    # ── Top row: incoming token stream ────────────────────────────────────
    add_text(slide, "Token stream",
             inner_x, y + Inches(0.10), inner_w, Inches(0.28),
             size=10, bold=True, color=THEME["secondary"])

    tokens = [
        ("سلام",    THEME["green"]),
        ("آپ",      THEME["green"]),
        ("بازر",    THEME["primary"]),   # OOV
        ("کیسے",    THEME["green"]),
        ("ہیں",     THEME["green"]),
    ]
    n = len(tokens)
    gap = Inches(0.08)
    tok_w = (inner_w - gap * (n - 1)) / n
    tok_h = Inches(0.50)
    tok_y = y + Inches(0.45)
    oov_cx = None
    cx = inner_x
    for txt, col in tokens:
        is_oov = (col == THEME["primary"])
        _chip(slide, txt, cx, tok_y, tok_w, tok_h, fill=col,
              size=14, bold=True)
        if is_oov:
            oov_cx = cx + tok_w / 2
            add_text(slide, "OOV", cx, tok_y - Inches(0.30),
                     tok_w, Inches(0.25),
                     size=9, bold=True, color=THEME["primary"],
                     align=PP_ALIGN.CENTER)
        cx += tok_w + gap

    # ── Vertical arrow into BK-tree ───────────────────────────────────────
    if oov_cx is not None:
        _line_between(slide, oov_cx, tok_y + tok_h,
                      oov_cx, tok_y + tok_h + Inches(0.4),
                      THEME["primary"], width=1.6)
        add_text(slide, "↓",
                 oov_cx - Inches(0.15),
                 tok_y + tok_h + Inches(0.05),
                 Inches(0.3), Inches(0.35),
                 size=18, bold=True, color=THEME["primary"],
                 align=PP_ALIGN.CENTER)

    # ── BK-tree hub + branching candidates ────────────────────────────────
    hub_y = tok_y + tok_h + Inches(0.55)
    hub_w = Inches(1.6); hub_h = Inches(0.85)
    hub_x = (oov_cx - hub_w / 2) if oov_cx else inner_x + (inner_w - hub_w) / 2

    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, hub_x, hub_y, hub_w, hub_h)
    hub.fill.solid(); hub.fill.fore_color.rgb = THEME["secondary"]
    hub.line.color.rgb = THEME["secondary"]; hub.line.width = Pt(1.5)
    hub.shadow.inherit = False
    add_text(slide, "BK-tree", hub_x, hub_y, hub_w, hub_h,
             size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, "~500K-word Urdu corpus",
             hub_x - Inches(0.7), hub_y + hub_h + Inches(0.05),
             hub_w + Inches(1.4), Inches(0.28),
             size=9, italic=True, color=THEME["text_dim"],
             align=PP_ALIGN.CENTER)

    # candidate cards on right
    cand_x = inner_x + Inches(0.05)
    cand_w = Inches(2.5)
    cand_y = hub_y + hub_h + Inches(0.55)
    cand_h = Inches(0.45)

    candidates = [
        ("بازار",   "edit=1 · top-1",  THEME["green"],  True),
        ("بازی",    "edit=1 · uni-hit", THEME["accent"], False),
        ("بزرگ",    "edit=2 · weak",   THEME["text_dim"], False),
    ]

    # Title bar
    add_text(slide, "Candidates · ranked (edit · uni · tri · bi · freq)",
             inner_x, cand_y - Inches(0.30), inner_w, Inches(0.28),
             size=10, bold=True, color=THEME["secondary"])

    # draw three candidate rows in a column on the right; lines from hub to each
    col_x = inner_x + inner_w - cand_w - Inches(0.1)
    yy = cand_y
    cand_centers = []
    for word, meta, col, top in candidates:
        box_h = cand_h
        box = add_card(slide, col_x, yy, cand_w, box_h,
                       fill=THEME["card"],
                       line=col, corner=0.20)
        box.line.width = Pt(1.6 if top else 0.8)
        add_text(slide, word, col_x + Inches(0.10), yy,
                 Inches(1.1), box_h,
                 size=13, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, meta, col_x + Inches(1.20), yy,
                 cand_w - Inches(1.30), box_h,
                 size=9, color=THEME["text_dim"], anchor=MSO_ANCHOR.MIDDLE,
                 font="Consolas")
        if top:
            add_text(slide, "★", col_x + cand_w - Inches(0.35), yy,
                     Inches(0.3), box_h,
                     size=14, bold=True, color=THEME["accent"],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cand_centers.append((col_x, yy + box_h / 2))
        yy += box_h + Inches(0.12)

    # connecting lines from hub to candidates
    for ccx, ccy in cand_centers:
        _line_between(slide,
                      hub_x + hub_w, hub_y + hub_h / 2,
                      ccx, ccy,
                      THEME["secondary"], width=0.9, dashed=True)

    # Selected output chip below
    sel_y = yy + Inches(0.1)
    add_text(slide, "Selected",
             inner_x, sel_y, Inches(1.5), Inches(0.3),
             size=10, bold=True, color=THEME["green"])
    sel_chip_x = inner_x + Inches(1.4)
    sel_chip_w = Inches(2.0)
    _chip(slide, "بازار", sel_chip_x, sel_y - Inches(0.05),
          sel_chip_w, Inches(0.45),
          fill=THEME["green"], size=14, bold=True)


# ============== Slide 1: Title ===============================================
def slide_title():
    s, bg = page(1)
    # Decorative blobs
    b1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(10.4), Inches(-1.5), Inches(4.5), Inches(4.5))
    b1.line.fill.background(); b1.fill.solid()
    b1.fill.fore_color.rgb = THEME["primary"]; b1.shadow.inherit=False
    _set_transparency(b1, 60)

    b2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(-2), Inches(5), Inches(4.5), Inches(4.5))
    b2.line.fill.background(); b2.fill.solid()
    b2.fill.fore_color.rgb = THEME["secondary"]; b2.shadow.inherit=False
    _set_transparency(b2, 65)

    b3 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(11.8), Inches(5.6), Inches(2.5), Inches(2.5))
    b3.line.fill.background(); b3.fill.solid()
    b3.fill.fore_color.rgb = THEME["accent"]; b3.shadow.inherit=False
    _set_transparency(b3, 75)

    # Brand pill
    pill = add_pill(s, "FAST-NUCES  •  FINAL YEAR PROJECT  •  2025-2026",
                    Inches(0.7), Inches(0.55), Inches(5.4), Inches(0.4),
                    THEME["primary"], RGBColor(0xFF, 0xFF, 0xFF), size=10)

    # Brand logo (if available)
    logo_path = FRONT_DIR / "coral-logo.png"
    if logo_path.exists():
        try:
            add_image(s, logo_path,
                      Inches(0.7), Inches(1.40), w=Inches(1.0), h=Inches(0.8))
        except Exception:
            pass

    # Big project tag
    add_text(s, "CORAL", Inches(1.9), Inches(1.2), Inches(11), Inches(1.4),
             size=96, bold=True, color=THEME["text"], font="Calibri")

    # Subtitle
    add_text(s,
        "Consensus-Based Refinement and Output Realignment\nfor Low-Resource Urdu ASR Post-Processing",
        Inches(0.7), Inches(2.7), Inches(12), Inches(1.4),
        size=24, color=THEME["text"], font="Calibri", line_spacing=1.2)

    # Tagline
    add_text(s,
        "A five-stage post-processing pipeline reducing Urdu ASR WER by up to 46.5% relative — without retraining any acoustic model.",
        Inches(0.7), Inches(4.25), Inches(11), Inches(1.0),
        size=15, color=THEME["text_dim"], italic=True, font="Calibri")

    # KPI strip
    kpis = [
        ("46.5%", "Rel. WER ↓\nConversational"),
        ("22.3%", "Rel. WER ↓\nCommon Voice"),
        ("500K", "Word Urdu\nCorpus"),
        ("5", "Pipeline\nStages"),
    ]
    x = Inches(0.7); w = Inches(2.7); gap = Inches(0.25)
    for v, lbl in kpis:
        c = add_card(s, x, Inches(5.4), w, Inches(1.25),
                     fill=THEME["card"], line=THEME["line"])
        add_text(s, v, x, Inches(5.45), w, Inches(0.6),
                 size=28, bold=True, color=THEME["primary"], align=PP_ALIGN.CENTER)
        add_text(s, lbl, x, Inches(6.05), w, Inches(0.6),
                 size=11, color=THEME["text"], align=PP_ALIGN.CENTER,
                 line_spacing=1.1)
        x = x + w + gap

    # Bottom band — team
    add_text(s, "Project Team",
             Inches(0.7), Inches(6.85), Inches(4), Inches(0.3),
             size=10, bold=True, color=THEME["secondary"])
    add_text(s, "Ali Irfan  21I-2572   •   Nouman Hafeez  21I-0416   •   Rafay Khattak  21I-0423",
             Inches(0.7), Inches(7.10), Inches(12), Inches(0.3),
             size=11, color=THEME["text"])

    set_slide_transition(s, "fade")
    inject_fade_animations(s,
        [int(pill.shape_id)] + collect_anim(s, exclude={int(bg.shape_id),
            int(b1.shape_id), int(b2.shape_id), int(b3.shape_id),
            int(pill.shape_id)}),
        delay_ms=120, base_ms=400)


# ============== Slide 2: Team & Supervisors ==================================
def slide_team():
    s, bg = page(2)
    add_decoration(s)
    add_header(s, "Meet the Team", "Three engineers. One Urdu ASR mission.")

    team = [
        ("Ali Irfan", "21I-2572",
         "Pipeline architect • Backend & deployment",
         "Stage 1 weighted split-merge alignment\nFastAPI orchestrator • HuggingFace Space deployment",
         THEME["primary"], "ali-irfan.jpg"),
        ("Nouman Hafeez", "21I-0416",
         "Corpus + retrieval • Web app • LLM strategy",
         "Stage 2 BK-tree + 500K n-gram corpus\nFull web app redesign • System architecture\nLLM research & integration strategy",
         THEME["secondary"], "nouman-hafeez.png"),
        ("Rafay Khattak", "21I-0423",
         "Research pitch • Evaluation lead",
         "Pitched the original CORAL research idea\nEvaluation methodology • C0→C7 ablation\nResidual-error analysis • WER/CER reporting",
         THEME["accent"], "rafay-khatak.jpg"),
    ]
    x = Inches(0.7); w = Inches(4.0); gap = Inches(0.2)
    for name, roll, role, work, color, photo in team:
        c = add_card(s, x, Inches(2.2), w, Inches(3.7),
                     fill=THEME["card"], line=THEME["line"], corner=0.06)
        # Top color band
        band = add_rect(s, x, Inches(2.2), w, Inches(0.55),
                        color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        band.adjustments[0] = 0.15

        # Avatar — real photo if available, else initials
        photo_path = FRONT_DIR / photo if photo else None
        if photo_path is not None and photo_path.exists():
            # backing circle for soft edge
            ring = add_rect(s, x + Inches(0.28), Inches(2.93),
                            Inches(0.99), Inches(0.99),
                            color, shape=MSO_SHAPE.OVAL)
            ring.line.fill.background()
            try:
                add_image(s, photo_path,
                          x + Inches(0.33), Inches(2.98),
                          w=Inches(0.89), h=Inches(0.89))
            except Exception:
                pass
        else:
            initials = "".join([p[0] for p in name.split()])
            add_rect(s, x + Inches(0.3), Inches(2.95),
                     Inches(0.95), Inches(0.95),
                     color, shape=MSO_SHAPE.OVAL)
            add_text(s, initials, x + Inches(0.3), Inches(2.95),
                     Inches(0.95), Inches(0.95), size=22, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

        # Name + roll
        add_text(s, name, x + Inches(1.4), Inches(3.0),
                 Inches(2.7), Inches(0.5), size=18, bold=True,
                 color=THEME["text"])
        add_text(s, roll, x + Inches(1.4), Inches(3.45),
                 Inches(2.7), Inches(0.3), size=11, color=THEME["text_dim"])
        # Role pill
        add_pill(s, role, x + Inches(0.3), Inches(4.1),
                 Inches(3.4), Inches(0.4), color,
                 RGBColor(0xFF, 0xFF, 0xFF), size=9)
        # Work
        add_text(s, work, x + Inches(0.3), Inches(4.65),
                 Inches(3.4), Inches(1.25), size=10.5,
                 color=THEME["text"], line_spacing=1.35)
        x = x + w + gap

    # Supervisor strip at bottom
    sup = add_card(s, Inches(0.7), Inches(6.05), Inches(12.0),
                   Inches(0.95), fill=THEME["card"], line=THEME["line"])
    add_text(s, "SUPERVISORS  •  ARTIFICIAL INTELLIGENCE & DATA SCIENCE",
             Inches(0.95), Inches(6.13), Inches(11), Inches(0.3),
             size=10, bold=True, color=THEME["secondary"])
    add_text(s, "Ms. Kainat Iqbal", Inches(0.95), Inches(6.38),
             Inches(4), Inches(0.4), size=15, bold=True, color=THEME["text"])
    add_text(s, "Supervisor", Inches(0.95), Inches(6.65), Inches(4),
             Inches(0.3), size=10, color=THEME["text_dim"], italic=True)
    add_text(s, "Ms. Saira Qamar", Inches(7), Inches(6.38),
             Inches(5), Inches(0.4), size=15, bold=True, color=THEME["text"])
    add_text(s, "Co-Supervisor", Inches(7), Inches(6.65), Inches(5),
             Inches(0.3), size=10, color=THEME["text_dim"], italic=True)

    add_footer(s, 2, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)})[:14],
                           delay_ms=180, base_ms=350)


# ============== Slide 3: Agenda ==============================================
def slide_agenda():
    s, bg = page(3)
    add_decoration(s)
    add_header(s, "Agenda", "What we will cover today")

    items = [
        ("01", "The Urdu ASR Problem", "Why does Urdu still hurt at 18-20% WER?"),
        ("02", "Research Hypothesis", "Five algorithmic levers, zero retraining."),
        ("03", "From FYP-1 to FYP-2", "What survived. What we dropped. And why."),
        ("04", "The CORAL Pipeline", "Stages 0-4 end-to-end."),
        ("05", "System Architecture", "Kaggle GPUs → FastAPI → Next.js."),
        ("06", "Empirical Results", "Common Voice + Conversational Urdu."),
        ("07", "Ablation & Residuals", "Which component buys which percentage?"),
        ("08", "Conclusion & Future Work", "Where CORAL goes next."),
    ]
    x = Inches(0.7); y = Inches(2.2)
    w = Inches(6.0); h = Inches(1.05); gap_y = Inches(0.15)
    col = 0
    for num, head, sub in items:
        cx = x if col == 0 else x + w + Inches(0.25)
        cy = y + (col * 0 + (items.index((num, head, sub)) // 2) * (h + gap_y))
        c = add_card(s, cx, cy, w, h, fill=THEME["card"],
                     line=THEME["line"], corner=0.10)
        # Number pill
        np_ = add_rect(s, cx + Inches(0.2), cy + Inches(0.2),
                       Inches(0.75), Inches(0.65),
                       THEME["primary"] if col == 0 else THEME["secondary"],
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        np_.adjustments[0] = 0.25
        add_text(s, num, cx + Inches(0.2), cy + Inches(0.2),
                 Inches(0.75), Inches(0.65), size=18, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # head + sub
        add_text(s, head, cx + Inches(1.1), cy + Inches(0.15),
                 w - Inches(1.2), Inches(0.42),
                 size=15, bold=True, color=THEME["text"])
        add_text(s, sub, cx + Inches(1.1), cy + Inches(0.57),
                 w - Inches(1.2), Inches(0.42),
                 size=10.5, color=THEME["text_dim"])
        col = (col + 1) % 2

    add_footer(s, 3, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 4: The Problem =========================================
def slide_problem():
    s, bg = page(4)
    add_decoration(s)
    add_header(s, "The Problem", "Urdu remains a severely under-resourced ASR language")

    # Left column: stats grid
    stats = [
        ("230 M+", "speakers worldwide", THEME["primary"]),
        ("13-20%", "WER for SOTA models", THEME["accent"]),
        ("0", "production-grade\nUrdu post-processors", THEME["secondary"]),
        ("1", "FYP team that took it on", THEME["violet"]),
    ]
    x = Inches(0.7); y = Inches(2.4); w = Inches(2.7); h = Inches(2.0)
    for i, (val, lbl, color) in enumerate(stats):
        cx = x + (i % 2) * (w + Inches(0.2))
        cy = y + (i // 2) * (h + Inches(0.2))
        add_card(s, cx, cy, w, h, fill=THEME["card"], line=THEME["line"])
        add_text(s, val, cx, cy + Inches(0.25), w, Inches(0.9),
                 size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, lbl, cx, cy + Inches(1.15), w, Inches(0.9),
                 size=12, color=THEME["text"], align=PP_ALIGN.CENTER,
                 line_spacing=1.2)

    # Right column: narrative
    right_x = Inches(6.7); right_w = Inches(6.0)
    add_text(s, "Foundation ASR models hit sub-5% WER on English.",
             right_x, Inches(2.4), right_w, Inches(0.5),
             size=15, bold=True, color=THEME["text"])
    add_text(s,
        "Yet on Urdu — the national language of Pakistan, written in Nastaliq script — "
        "Whisper, MMS and SeamlessM4T still hover between 13% (read speech) and 20% (conversational). "
        "Critical domains — healthcare, legal, accessibility, education — remain blocked.",
        right_x, Inches(2.85), right_w, Inches(2.2),
        size=13, color=THEME["text"], line_spacing=1.35)

    add_text(s, "Why is Urdu hard?",
             right_x, Inches(5.05), right_w, Inches(0.4),
             size=14, bold=True, color=THEME["secondary"])
    challenges = [
        "• Limited training data — Urdu is a minority of multilingual mixtures",
        "• Arabic ↔ Urdu Unicode conflation (U+0643 vs U+06A9, etc.)",
        "• Word-boundary disagreement across architectures",
        "• Rampant Urdu ↔ English code-switching in conversation",
    ]
    add_text(s, "\n".join(challenges),
             right_x, Inches(5.4), right_w, Inches(1.6),
             size=11.5, color=THEME["text"], line_spacing=1.45)

    add_footer(s, 4, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 5: Hypothesis ==========================================
def slide_hypothesis():
    s, bg = page(5)
    # Center decorative band
    band = add_rect(s, 0, Inches(2.4), SW, Inches(3.0),
                    THEME["card"], shape=MSO_SHAPE.RECTANGLE)
    add_decoration(s)
    add_header(s, "Research Hypothesis",
               "Five algorithmic levers — no fine-tuning needed")

    # Big quote
    add_text(s, "“",
             Inches(0.6), Inches(2.3), Inches(1), Inches(1),
             size=96, bold=True, color=THEME["primary"])
    add_text(s,
        "Combining (i) Urdu-specific Unicode normalisation, "
        "(ii) weighted split-merge-aware alignment, "
        "(iii) hybrid OOV detection with BK-tree fuzzy lookup ranked by an n-gram LM, "
        "(iv) conservative position-wise voting, and "
        "(v) targeted LLM refinement of the voted output —\n\n"
        "is sufficient to substantially reduce WER over the strongest single-model baseline, "
        "on both read-speech and conversational Urdu, without any fine-tuning of the acoustic models.",
        Inches(1.6), Inches(2.6), Inches(11.0), Inches(2.7),
        size=17, italic=True, color=THEME["text"], line_spacing=1.4)

    # five-lever footer
    levers = [("0", "Normalize", THEME["primary"]),
              ("1", "Align", THEME["secondary"]),
              ("2", "OOV+BK-tree", THEME["accent"]),
              ("3", "Vote", THEME["green"]),
              ("4", "LLM Refine", THEME["violet"])]
    x = Inches(1.0); w = Inches(2.25); gap = Inches(0.15)
    for n, label, color in levers:
        c = add_card(s, x, Inches(5.7), w, Inches(0.95),
                     fill=THEME["card"], line=color, corner=0.12)
        c.line.width = Pt(2.0)
        add_text(s, n, x + Inches(0.15), Inches(5.78), Inches(0.7),
                 Inches(0.8), size=32, bold=True, color=color)
        add_text(s, label, x + Inches(0.9), Inches(5.9), w - Inches(1.0),
                 Inches(0.6), size=13, bold=True, color=THEME["text"],
                 anchor=MSO_ANCHOR.MIDDLE)
        x = x + w + gap

    add_footer(s, 5, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id),
                                                        int(band.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 6: FYP1 vs FYP2 Evolution ==============================
def slide_evolution():
    s, bg = page(6)
    add_decoration(s)
    add_header(s, "Evolution", "From FYP-1 Generate-and-Refine  →  FYP-2 five-stage pipeline")

    # Two columns
    left_x = Inches(0.7); right_x = Inches(7.0); w = Inches(5.9); y = Inches(2.2)
    add_card(s, left_x, y, w, Inches(4.5), fill=THEME["card"],
             line=THEME["accent"], corner=0.05)
    add_card(s, right_x, y, w, Inches(4.5), fill=THEME["card"],
             line=THEME["primary"], corner=0.05)

    # Left header
    add_pill(s, "FYP-1  •  GENERATE-AND-REFINE", left_x + Inches(0.3),
             y + Inches(0.25), Inches(3.8), Inches(0.4),
             THEME["accent"], RGBColor(0x1F, 0x29, 0x37), size=11)
    add_text(s, "2-stage architecture",
             left_x + Inches(0.3), y + Inches(0.8), w, Inches(0.4),
             size=15, bold=True, color=THEME["text"])
    fyp1_items = [
        ("Stage 1", "Confidence-annotated multi-model ASR"),
        ("Stage 2", "Alif-1.0-8B LLM arbitrates between hypotheses"),
        ("Eval", "20-sample synthetic dataset, 3 calibration regimes"),
        ("Metric", "WER + ECE (Expected Calibration Error)"),
        ("Result", "Whisper-Large baseline 17.76% WER · LLM 18.15%"),
    ]
    yy = y + Inches(1.4)
    for k, v in fyp1_items:
        add_pill(s, k, left_x + Inches(0.3), yy, Inches(0.9),
                 Inches(0.32), THEME["bg"], THEME["accent"], size=9.5)
        add_text(s, v, left_x + Inches(1.35), yy - Inches(0.02),
                 w - Inches(1.5), Inches(0.5), size=11.5,
                 color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + Inches(0.55)

    # Right header
    add_pill(s, "FYP-2  •  FIVE-STAGE POST-PROCESSING",
             right_x + Inches(0.3), y + Inches(0.25), Inches(4.6),
             Inches(0.4), THEME["primary"], RGBColor(0xFF, 0xFF, 0xFF), size=11)
    add_text(s, "Compositional algorithmic pipeline",
             right_x + Inches(0.3), y + Inches(0.8), w, Inches(0.4),
             size=15, bold=True, color=THEME["text"])
    fyp2_items = [
        ("Stage 0", "Urdu normalisation (Unicode, diacritics, tatweel)"),
        ("Stage 1", "Weighted split-merge alignment (SAME/SPLIT/MERGE/NOISE)"),
        ("Stage 2", "BK-tree OOV + n-gram ranking"),
        ("Stage 3", "Position-wise voting, conservative tie-break"),
        ("Stage 4", "Frontier LLM refinement (Gemini / Claude / GPT)"),
    ]
    yy = y + Inches(1.4)
    for k, v in fyp2_items:
        add_pill(s, k, right_x + Inches(0.3), yy, Inches(0.9),
                 Inches(0.32), THEME["bg"], THEME["primary"], size=9.5)
        add_text(s, v, right_x + Inches(1.35), yy - Inches(0.02),
                 w - Inches(1.5), Inches(0.5), size=11.5,
                 color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + Inches(0.55)

    # Arrow connector
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.55),
                              Inches(4.3), Inches(0.45), Inches(0.35))
    arr.fill.solid(); arr.fill.fore_color.rgb = THEME["primary"]
    arr.line.fill.background(); arr.shadow.inherit = False

    add_footer(s, 6, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)})[:14],
                           delay_ms=140, base_ms=300)


# ============== Slide 7: Components Dropped ==================================
def slide_dropped():
    s, bg = page(7)
    add_decoration(s)
    add_header(s, "What We Dropped",
               "Five FYP-1 elements removed or repositioned for FYP-2")

    items = [
        ("Word-level confidence extraction",
         "Whisper/Wav2Vec2 raw confidences live on incompatible scales — proved unreliable in practice.",
         "Replaced with static, model-quality based weighting only at the LLM stage.",
         THEME["primary"]),
        ("Expected Calibration Error (ECE)",
         "Useful for FYP-1 confidence-quality comparison, but tangential to deployed pipeline.",
         "FYP-2 reports WER + CER only. ECE retained as historical context.",
         THEME["accent"]),
        ("Synthetic 20-sample evaluation",
         "Three calibration regimes × three error types ignored Urdu-specific failures.",
         "Replaced with 2 995 Common Voice + 500 Conversational Urdu clips.",
         THEME["secondary"]),
        ("Alif-1.0-8B-Instruct LLM",
         "Heavy GPU + weaker instruction following than commercial frontier models.",
         "Replaced with client-side Gemini / OpenRouter (Claude, GPT-4o).",
         THEME["violet"]),
        ("Two-stage Generate-and-Refine framing",
         "Treated LLM as primary fusion mechanism — wrong abstraction for Urdu.",
         "Reframed as 5 sequential stages with LLM as the LAST refinement step.",
         THEME["green"]),
    ]
    x = Inches(0.7); w = Inches(12.0); h = Inches(0.85); y = Inches(2.15); gap = Inches(0.10)
    for title, was, now, color in items:
        c = add_card(s, x, y, w, h, fill=THEME["card"],
                     line=THEME["line"], corner=0.10)
        # left color tag
        tag = add_rect(s, x, y, Inches(0.18), h, color,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tag.adjustments[0] = 0.4
        # strikethrough title
        tb = s.shapes.add_textbox(x + Inches(0.4), y + Inches(0.05),
                                  Inches(4.6), Inches(0.45))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.name = "Calibri"; r.font.size = Pt(13.5); r.font.bold = True
        r.font.color.rgb = THEME["text"]
        # strike-through via XML
        rPr = r._r.get_or_add_rPr()
        rPr.set("strike", "sngStrike")
        # then
        add_text(s, f"× was:  {was}",
                 x + Inches(0.4), y + Inches(0.45),
                 Inches(6.2), Inches(0.35),
                 size=10, color=THEME["text_dim"])
        add_text(s, f"→ now:  {now}",
                 x + Inches(6.6), y + Inches(0.07),
                 Inches(5.2), Inches(0.7),
                 size=11, color=color, bold=True, line_spacing=1.2)
        y = y + h + gap

    add_footer(s, 7, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)})[:14],
                           delay_ms=140, base_ms=300)


# ============== Slide 8: Pipeline overview ====================================
def slide_pipeline():
    s, bg = page(8)
    add_decoration(s)
    add_header(s, "The CORAL Pipeline", "Five stages, end-to-end")

    stages = [
        ("0", "Normalize", "Urdu Unicode\ndiacritics, tatweel,\nASCII strip",
         THEME["primary"]),
        ("1", "Align", "Weighted split-merge\nLevenshtein\nSAME/SPLIT/MERGE/NOISE",
         THEME["secondary"]),
        ("2", "Detect OOV", "BK-tree fuzzy +\nbi/tri-gram rank\non 500K vocab",
         THEME["accent"]),
        ("3", "Vote", "Position-wise majority\nconservative tie-break\nOOV substitution",
         THEME["green"]),
        ("4", "LLM Refine", "Fluency, code-switch,\nnamed-entity recovery\n(Gemini/Claude/GPT)",
         THEME["violet"]),
    ]
    x = Inches(0.5); w = Inches(2.45); gap = Inches(0.12)
    for n, head, body, color in stages:
        c = add_card(s, x, Inches(2.5), w, Inches(3.6),
                     fill=THEME["card"], line=color, corner=0.08)
        c.line.width = Pt(2.5)
        # circle stage number
        circ = add_rect(s, x + Inches(0.85), Inches(2.65), Inches(0.75),
                        Inches(0.75), color, shape=MSO_SHAPE.OVAL)
        add_text(s, n, x + Inches(0.85), Inches(2.65),
                 Inches(0.75), Inches(0.75), size=22, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, head, x, Inches(3.5), w, Inches(0.5),
                 size=17, bold=True, color=THEME["text"], align=PP_ALIGN.CENTER)
        add_text(s, body, x + Inches(0.15), Inches(4.05),
                 w - Inches(0.3), Inches(1.9), size=11,
                 color=THEME["text_dim"], align=PP_ALIGN.CENTER,
                 line_spacing=1.3)
        if n != "4":
            # arrow
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + w - Inches(0.0),
                                      Inches(4.0), Inches(0.12),
                                      Inches(0.35))
            arr.fill.solid(); arr.fill.fore_color.rgb = THEME["text_dim"]
            arr.line.fill.background(); arr.shadow.inherit = False
        x = x + w + gap

    # Input / Output bands
    add_pill(s, "INPUT  •  raw multi-model ASR hypotheses",
             Inches(0.5), Inches(2.05), Inches(5.0), Inches(0.32),
             THEME["bg"], THEME["secondary"], size=11)
    add_pill(s, "OUTPUT  •  corrected Urdu transcript",
             Inches(7.8), Inches(2.05), Inches(5.0), Inches(0.32),
             THEME["bg"], THEME["primary"], size=11)

    # Bottom band: key numbers
    nums = add_card(s, Inches(0.5), Inches(6.3), Inches(12.3),
                    Inches(0.75), fill=THEME["card"], line=THEME["line"])
    add_text(s,
        "  Server-side  Stages 0-3  •   Browser-side  Stage 4  •   "
        "Audio  →  /transcribe  →  /align  →  /oov  →  /correct  →  LLM",
        Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55),
        size=12.5, color=THEME["text"], align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 8, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=160, base_ms=300)


# ============== Slide 9: Stage 0 — Normalisation =============================
def slide_stage0():
    s, bg = page(9)
    add_decoration(s)
    add_header(s, "Stage 0  •  Normalisation",
               "Strip the noise — make Urdu look like Urdu")

    # left: bullets
    bullets = [
        ("Diacritics", "Strip U+0610–U+061A, U+064B–U+065F, U+0670 (superscript alef)"),
        ("Arabic → Urdu", "7 canonical character maps: ك→ک, ي→ی, ة→ت, ﻻ→لا …"),
        ("Tatweel", "Remove U+0640 elongation character"),
        ("Punctuation", "Drop Arabic + Urdu + ASCII punctuation + ZWJ/ZWNJ"),
        ("Whitespace", "Collapse to single spaces · NFC normalise"),
        ("Effect", "1.9 WER points  ↓  on Conversational Urdu — for free."),
    ]
    x = Inches(0.7); y = Inches(2.4); w = Inches(6.0)
    for k, v in bullets:
        c = add_card(s, x, y, w, Inches(0.6), fill=THEME["card"],
                     line=THEME["line"], corner=0.18)
        add_pill(s, k, x + Inches(0.15), y + Inches(0.12), Inches(1.3),
                 Inches(0.36), THEME["primary"],
                 RGBColor(0xFF, 0xFF, 0xFF), size=10)
        add_text(s, v, x + Inches(1.6), y + Inches(0.12),
                 w - Inches(1.7), Inches(0.45),
                 size=12, color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.7)

    # right: code-style example
    code_card = add_card(s, Inches(7.0), Inches(2.4), Inches(5.7),
                         Inches(4.4), fill=RGBColor(0x12, 0x1A, 0x28),
                         line=THEME["line"], corner=0.05)
    add_text(s, "EXAMPLE  /  before  →  after",
             Inches(7.2), Inches(2.5), Inches(5), Inches(0.3),
             size=10, bold=True, color=THEME["secondary"])
    code_before = "Raw:    سلامؐ علیکم،   ﻛیسے ہیں — ٹھیک؟"
    code_after  = "Norm:  سلام علیکم  کیسے ہیں  ٹھیک"
    add_text(s, code_before, Inches(7.2), Inches(2.95), Inches(5.4),
             Inches(0.45), size=14, font="Consolas",
             color=THEME["accent"])
    add_text(s, "↓",
             Inches(7.2), Inches(3.4), Inches(5.4), Inches(0.4),
             size=18, bold=True, color=THEME["primary"], align=PP_ALIGN.CENTER)
    add_text(s, code_after, Inches(7.2), Inches(3.85), Inches(5.4),
             Inches(0.45), size=14, font="Consolas", color=THEME["secondary"])

    # quick stats
    add_text(s, "Why Stage 0 matters",
             Inches(7.2), Inches(4.5), Inches(5.4), Inches(0.35),
             size=13, bold=True, color=THEME["text"])
    add_text(s,
        "A large fraction of apparent ASR errors are NOT recognition failures —\n"
        "they are Unicode-variant mismatches against the reference. Stage 0 dissolves\n"
        "this measurement artefact before any downstream voting decision.",
        Inches(7.2), Inches(4.85), Inches(5.4), Inches(1.9),
        size=11.5, color=THEME["text_dim"], line_spacing=1.4)

    add_footer(s, 9, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 10: Stage 1 — Split-Merge Alignment =====================
def slide_stage1():
    s, bg = page(10)
    add_decoration(s)
    add_header(s, "Stage 1  •  Split-Merge Alignment",
               "Word boundaries are first-class citizens")

    # Left: textual description
    left_x = Inches(0.7); left_w = Inches(6.0)
    add_text(s, "The Core Insight",
             left_x, Inches(2.3), left_w, Inches(0.4),
             size=15, bold=True, color=THEME["secondary"])
    add_text(s,
        "Different ASR models disagree on WHERE word boundaries fall. "
        "One model emits a phrase as a single token; another splits it in two; "
        "a third merges across syllables. ROVER mis-reads this as substitution.",
        left_x, Inches(2.7), left_w, Inches(1.6),
        size=12.5, color=THEME["text"], line_spacing=1.4)

    add_text(s, "Per-chunk metadata produced",
             left_x, Inches(4.4), left_w, Inches(0.4),
             size=14, bold=True, color=THEME["text"])
    tags = [
        ("SAME",  "tokens align 1:1",                THEME["green"]),
        ("SPLIT", "one source token → many target",  THEME["accent"]),
        ("MERGE", "many source → one target",        THEME["orange"]),
        ("NOISE", "neither side aligns reliably",    THEME["primary"]),
    ]
    yy = Inches(4.85)
    for t, d, color in tags:
        add_pill(s, t, left_x, yy, Inches(1.0), Inches(0.35),
                 color, RGBColor(0xFF, 0xFF, 0xFF), size=10.5)
        add_text(s, d, left_x + Inches(1.15), yy - Inches(0.02),
                 left_w - Inches(1.2), Inches(0.4),
                 size=12, color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + Inches(0.43)

    add_text(s,
        "→ 36.5% of all alignment events on Conversational Urdu are SPLIT or MERGE.\n"
        "→ Existing fusion methods (ROVER) lose this signal entirely.",
        left_x, Inches(6.65), left_w, Inches(0.9),
        size=11.5, italic=True, color=THEME["primary"], line_spacing=1.4)

    # Right: hand-drawn split-merge alignment diagram
    draw_splitmerge_diagram(s,
                            x=Inches(7.0), y=Inches(2.25),
                            w=Inches(5.85), h=Inches(4.55))

    add_footer(s, 10, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 11: Stage 2 — OOV + BK-tree =============================
def slide_stage2():
    s, bg = page(11)
    add_decoration(s)
    add_header(s, "Stage 2  •  OOV Detection + BK-tree",
               "Find the rare words. Correct them with context.")

    # Right: hand-drawn OOV + BK-tree diagram
    draw_oov_bktree_diagram(s,
                            x=Inches(7.0), y=Inches(2.25),
                            w=Inches(5.85), h=Inches(4.55))

    # Left content
    left_x = Inches(0.7); left_w = Inches(6.0)
    add_text(s, "OOV criterion (FR3.1)",
             left_x, Inches(2.3), left_w, Inches(0.4),
             size=14, bold=True, color=THEME["secondary"])
    add_text(s,
        "Flag a token as OOV iff it appears in EXACTLY ONE model's output\n"
        "AND its corpus frequency is below 2 000.",
        left_x, Inches(2.7), left_w, Inches(0.9),
        size=12, color=THEME["text"], line_spacing=1.4)

    add_text(s, "BK-tree fuzzy lookup",
             left_x, Inches(3.7), left_w, Inches(0.4),
             size=14, bold=True, color=THEME["secondary"])
    add_text(s,
        "Burkhard-Keller tree over a ~500K-word Urdu vocabulary corpus.\n"
        "Length-adaptive edit-distance threshold. >10× faster than linear scan.",
        left_x, Inches(4.1), left_w, Inches(1.0),
        size=12, color=THEME["text"], line_spacing=1.4)

    add_text(s, "Ranking signals (lexicographic)",
             left_x, Inches(5.2), left_w, Inches(0.4),
             size=14, bold=True, color=THEME["secondary"])
    rank = [
        "1.  Edit distance",
        "2.  Unigram presence",
        "3.  Trigram-context hits",
        "4.  Bigram-context hits",
        "5.  Corpus frequency",
    ]
    add_text(s, "\n".join(rank),
             left_x, Inches(5.6), left_w, Inches(1.8),
             size=12, color=THEME["text"], line_spacing=1.4)

    add_footer(s, 11, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 12: Stage 3 — Voting ===================================
def slide_stage3():
    s, bg = page(12)
    add_decoration(s)
    add_header(s, "Stage 3  •  Consensus Voting",
               "Position-wise majority — but conservative")

    bullets = [
        ("For each non-OOV source position",
         "Collect votes from companion models. Skip insertion-aligned slots."),
        ("Plurality wins — uniquely",
         "Replace the source word only if there is a UNIQUE plurality winner."),
        ("Tie-break = preserve",
         "Conservative: on ties, retain the source word. No silent substitutions."),
        ("For each OOV position",
         "Replace the source word with the top-ranked BK-tree candidate."),
        ("Output",
         "Corrected sentence + source + position-indexed diff."),
    ]
    x = Inches(0.7); y = Inches(2.4); w = Inches(8.0)
    for k, v in bullets:
        c = add_card(s, x, y, w, Inches(0.7), fill=THEME["card"],
                     line=THEME["line"], corner=0.10)
        add_text(s, "▸",
                 x + Inches(0.15), y + Inches(0.12), Inches(0.4),
                 Inches(0.4), size=18, bold=True, color=THEME["green"])
        add_text(s, k, x + Inches(0.55), y + Inches(0.10),
                 w - Inches(0.6), Inches(0.36), size=13, bold=True,
                 color=THEME["text"])
        add_text(s, v, x + Inches(0.55), y + Inches(0.42),
                 w - Inches(0.6), Inches(0.32), size=11,
                 color=THEME["text_dim"])
        y = y + Inches(0.8)

    # Right card: example
    add_card(s, Inches(9.0), Inches(2.4), Inches(3.7), Inches(4.4),
             fill=THEME["card"], line=THEME["green"], corner=0.05)
    add_text(s, "VOTING EXAMPLE",
             Inches(9.2), Inches(2.5), Inches(3.5), Inches(0.3),
             size=10, bold=True, color=THEME["green"])
    add_text(s, "Position 4",
             Inches(9.2), Inches(2.85), Inches(3.5), Inches(0.35),
             size=14, bold=True, color=THEME["text"])
    votes = [
        ("Seamless", "کتاب", THEME["secondary"]),
        ("Whisper-L", "کتاب", THEME["secondary"]),
        ("Whisper-M", "ﮐﺗﺎب", THEME["accent"]),
        ("Wav2Vec",  "کتاب", THEME["secondary"]),
    ]
    yy = Inches(3.3)
    for src, w_, color in votes:
        add_pill(s, src, Inches(9.2), yy, Inches(1.4),
                 Inches(0.32), color, RGBColor(0xFF, 0xFF, 0xFF), size=10)
        add_text(s, w_, Inches(10.7), yy - Inches(0.05), Inches(2.0),
                 Inches(0.45), size=15, color=THEME["text"],
                 font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + Inches(0.5)

    add_text(s, "Winner:  کتاب  (4 / 4)",
             Inches(9.2), Inches(5.6), Inches(3.5), Inches(0.5),
             size=14, bold=True, color=THEME["green"])

    add_footer(s, 12, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 13: Stage 4 — LLM Refinement ===========================
def slide_stage4():
    s, bg = page(13)
    add_decoration(s)
    add_header(s, "Stage 4  •  LLM Refinement",
               "Fluency, code-switching, named entities — handled last, in-browser")

    # Left: bullets
    left_x = Inches(0.7); left_w = Inches(6.2)
    add_text(s, "What the LLM handles",
             left_x, Inches(2.3), left_w, Inches(0.4),
             size=15, bold=True, color=THEME["violet"])
    pts = [
        "• Restore sentence-level fluency lost during voting",
        "• Recover code-switched English tokens stripped by Stage 0",
        "• Resolve grammar agreement that spans many words",
        "• Disambiguate named entities & proper nouns",
        "• Operate under explicit authority limits — can NOT add new content",
    ]
    add_text(s, "\n".join(pts),
             left_x, Inches(2.75), left_w, Inches(2.2),
             size=12.5, color=THEME["text"], line_spacing=1.5)

    add_text(s, "Why client-side?",
             left_x, Inches(5.0), left_w, Inches(0.4),
             size=14, bold=True, color=THEME["violet"])
    add_text(s,
        "• User-supplied API keys never leave the browser\n"
        "• No GPU cost on our backend\n"
        "• Easy provider switch — Gemini / OpenRouter / Claude / GPT-4o",
        left_x, Inches(5.4), left_w, Inches(1.5),
        size=12, color=THEME["text_dim"], line_spacing=1.5)

    # Right: provider cards
    providers = [
        ("Google", "Gemini 1.5 Pro\nGemini 2.0 Flash", THEME["blue"]),
        ("Anthropic", "Claude 3.5 Sonnet\nClaude 3 Opus",    THEME["orange"]),
        ("OpenAI",   "GPT-4o\nGPT-4o-mini",                  THEME["green"]),
        ("Meta",     "Llama 3.1 405B\nvia OpenRouter",       THEME["primary"]),
    ]
    rx = Inches(7.3); rw = Inches(2.7); rh = Inches(1.85); rgap = Inches(0.15)
    for i, (vendor, model, color) in enumerate(providers):
        cx = rx + (i % 2) * (rw + rgap)
        cy = Inches(2.3) + (i // 2) * (rh + rgap)
        c = add_card(s, cx, cy, rw, rh, fill=THEME["card"],
                     line=color, corner=0.08)
        c.line.width = Pt(2.0)
        add_pill(s, vendor, cx + Inches(0.2), cy + Inches(0.2),
                 rw - Inches(0.4), Inches(0.35), color,
                 RGBColor(0xFF, 0xFF, 0xFF), size=11)
        add_text(s, model, cx + Inches(0.2), cy + Inches(0.7),
                 rw - Inches(0.4), Inches(1.1),
                 size=11.5, color=THEME["text"], align=PP_ALIGN.CENTER,
                 line_spacing=1.3)

    # Bottom: prompt structure
    add_card(s, Inches(7.3), Inches(6.0), Inches(5.55), Inches(0.85),
             fill=RGBColor(0x12, 0x1A, 0x28), line=THEME["violet"])
    add_text(s, "PROMPT  =  voted transcript  +  align metadata  +  OOV candidates",
             Inches(7.3), Inches(6.1), Inches(5.55), Inches(0.7),
             size=11, color=THEME["violet"], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font="Consolas", bold=True)

    add_footer(s, 13, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 14: System Architecture ================================
def slide_architecture():
    s, bg = page(14)
    add_decoration(s)
    add_header(s, "System Architecture",
               "Four tiers, one rendezvous registry")

    img = add_image(s, FIG_DIR / "report_fig_4_1_architecture.png",
                    Inches(0.7), Inches(2.3), w=Inches(8.0))

    # Right legend / tiers
    rx = Inches(9.0); rw = Inches(3.7)
    tiers = [
        ("Inference Tier", "Kaggle GPU notebooks\n+ ngrok HTTPS tunnels", THEME["green"]),
        ("Backend Tier",   "FastAPI · /align /oov /correct\n+ Model registry", THEME["blue"]),
        ("Frontend Tier",  "Next.js · React · 4-pass UI\nClient-side LLM call", THEME["accent"]),
        ("Data Tier",      "HuggingFace Hub\n+ DuckDB n-gram tables", THEME["secondary"]),
    ]
    y = Inches(2.2)
    for k, v, c in tiers:
        card = add_card(s, rx, y, rw, Inches(1.05),
                        fill=THEME["card"], line=c, corner=0.08)
        card.line.width = Pt(2.0)
        add_pill(s, k, rx + Inches(0.15), y + Inches(0.15),
                 Inches(1.7), Inches(0.32), c,
                 RGBColor(0xFF, 0xFF, 0xFF), size=11)
        add_text(s, v, rx + Inches(0.15), y + Inches(0.5),
                 rw - Inches(0.3), Inches(0.6), size=10.5,
                 color=THEME["text"], line_spacing=1.3)
        y = y + Inches(1.15)

    add_footer(s, 14, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 15: Deployment topology ================================
def slide_deployment():
    s, bg = page(15)
    add_decoration(s)
    add_header(s, "System Topology",
               "Self-registration · heartbeats · async fan-out")

    # Official topology figure from the FYP2 report (Figure 4.2)
    img = add_image(s, FIG_DIR / "report_fig_4_2_topology.png",
                    Inches(0.55), Inches(2.25), w=Inches(9.0))

    # Right-side annotations
    rx = Inches(9.8); rw = Inches(3.0); ry = Inches(2.3); gap = Inches(0.15)
    notes = [
        ("Self-registration",
         "Kaggle notebooks POST /registry/register on boot · ping every 60 s",
         THEME["green"]),
        ("Fan-out",
         "POST /registry/transcribe dispatches audio to every live model in parallel · httpx.AsyncClient · 60 s timeout",
         THEME["secondary"]),
        ("Eviction",
         "Registry sweep every 20 s drops sessions silent for >300 s",
         THEME["primary"]),
        ("Auth",
         "X-Registry-Token for register/ping · X-API-Token on transcribe",
         THEME["violet"]),
    ]
    for k, v, color in notes:
        c = add_card(s, rx, ry, rw, Inches(1.05),
                     fill=THEME["card"], line=color, corner=0.08)
        c.line.width = Pt(1.5)
        add_pill(s, k, rx + Inches(0.15), ry + Inches(0.12),
                 rw - Inches(0.3), Inches(0.30),
                 color, RGBColor(0xFF, 0xFF, 0xFF), size=10)
        add_text(s, v, rx + Inches(0.15), ry + Inches(0.46),
                 rw - Inches(0.3), Inches(0.55),
                 size=9.5, color=THEME["text"], line_spacing=1.25)
        ry = ry + Inches(1.05) + gap

    # Bottom caption
    add_text(s,
        "Each ASR model runs on its own Kaggle GPU notebook behind a private ngrok HTTPS tunnel; the backend registry is the rendezvous point.  "
        "The browser talks only to the backend (and, for Stage 4, directly to the LLM provider).",
        Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.45),
        size=10.5, italic=True, color=THEME["text_dim"], line_spacing=1.3,
        align=PP_ALIGN.LEFT)

    add_footer(s, 15, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=130, base_ms=300)


# ============== Slide 16: Tech Stack =========================================
def slide_techstack():
    s, bg = page(16)
    add_decoration(s)
    add_header(s, "Technology Stack", "Modern, lean, production-ready")

    cols = [
        ("Backend", THEME["blue"],
         ["FastAPI 0.109 · uvicorn 0.23", "rapidfuzz 3.6 (Levenshtein)",
          "joblib 1.3 (BK-tree)", "duckdb 0.10 + pyarrow 14",
          "huggingface_hub 0.20", "httpx 0.26 (async)",
          "HuggingFace Spaces deploy", "Docker (port 7860)"]),
        ("Frontend", THEME["accent"],
         ["Next.js 16 · App Router", "React 19 · TypeScript 5",
          "TailwindCSS v4", "Multi-page R&D site",
          "Dark cinematic theme · aurora", "Animated waveform · particles",
          "MediaRecorder API · localStorage", "GPT-OSS chat endpoint"]),
        ("Inference", THEME["green"],
         ["transformers (HF) · torch", "librosa · soundfile",
          "Kaggle GPU notebooks (T4)", "ngrok HTTPS tunnels",
          "Whisper-Large-v3", "SeamlessM4T-v2-Large",
          "Wav2Vec2-Urdu", "Model registry auto-register"]),
        ("Evaluation", THEME["violet"],
         ["jiwer (WER/CER)", "Common Voice Urdu (2,995)",
          "Conversational Urdu (500)", "8-step ablation (C0→C7)",
          "Residual error annotation", "matplotlib + seaborn",
          "LaTeX + TikZ (this report)", "Reproducible eval TSV"]),
    ]
    x = Inches(0.5); w = Inches(3.05); gap = Inches(0.15)
    for head, color, items in cols:
        c = add_card(s, x, Inches(2.3), w, Inches(4.7),
                     fill=THEME["card"], line=color, corner=0.07)
        c.line.width = Pt(2.0)
        add_pill(s, head, x + Inches(0.2), Inches(2.45),
                 w - Inches(0.4), Inches(0.4),
                 color, RGBColor(0xFF, 0xFF, 0xFF), size=12)
        body = "\n".join([f"•  {i}" for i in items])
        add_text(s, body, x + Inches(0.25), Inches(2.95),
                 w - Inches(0.5), Inches(3.7),
                 size=11, color=THEME["text"], line_spacing=1.5)
        x = x + w + gap

    add_footer(s, 16, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 17: Evaluation Setup ===================================
def slide_eval_setup():
    s, bg = page(17)
    add_decoration(s)
    add_header(s, "Evaluation Setup",
               "Two benchmarks, four configurations, every component")

    # Benchmarks
    cards = [
        ("Common Voice Urdu",     "2,995 read-speech utterances",
         "Public, Mozilla-curated, read passages.",      THEME["secondary"]),
        ("Conversational Urdu",   "500 conversational clips",
         "Spontaneous speech with code-switching.",      THEME["primary"]),
    ]
    x = Inches(0.7); w = Inches(6.0); gap = Inches(0.3)
    for head, sub, body, color in cards:
        c = add_card(s, x, Inches(2.3), w, Inches(1.7),
                     fill=THEME["card"], line=color, corner=0.06)
        c.line.width = Pt(2.0)
        add_text(s, head, x + Inches(0.3), Inches(2.45),
                 w - Inches(0.6), Inches(0.45), size=18, bold=True,
                 color=THEME["text"])
        add_text(s, sub, x + Inches(0.3), Inches(2.9),
                 w - Inches(0.6), Inches(0.4), size=14,
                 color=color, bold=True)
        add_text(s, body, x + Inches(0.3), Inches(3.3),
                 w - Inches(0.6), Inches(0.6),
                 size=11.5, color=THEME["text_dim"], line_spacing=1.4)
        x = x + w + gap

    # Configurations
    add_text(s, "Pipeline configurations evaluated",
             Inches(0.7), Inches(4.4), Inches(12), Inches(0.4),
             size=14, bold=True, color=THEME["secondary"])
    configs = [
        ("untouched",     "Raw ASR output, no post-processing",        THEME["text_dim"]),
        ("norm_baseline", "Stage 0 normalisation only",               THEME["accent"]),
        ("robust",        "Stages 0–3 (normalise + align + OOV + vote)", THEME["primary"]),
        ("rrf",           "Reciprocal-rank-fusion variant",            THEME["violet"]),
    ]
    x = Inches(0.7); w = Inches(2.95); gap = Inches(0.1)
    for k, v, color in configs:
        c = add_card(s, x, Inches(4.9), w, Inches(1.9),
                     fill=THEME["card"], line=color, corner=0.07)
        c.line.width = Pt(1.5)
        add_pill(s, k, x + Inches(0.2), Inches(5.05), w - Inches(0.4),
                 Inches(0.35), color, RGBColor(0xFF, 0xFF, 0xFF), size=11)
        add_text(s, v, x + Inches(0.2), Inches(5.5), w - Inches(0.4),
                 Inches(1.2), size=11, color=THEME["text"], line_spacing=1.4)
        x = x + w + gap

    add_footer(s, 17, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 18: Read Speech Results ================================
def slide_results_readspeech():
    s, bg = page(18)
    add_decoration(s)
    add_header(s, "Results · Common Voice (read speech)",
               "WER reductions across every model")

    img = add_image(s, FIG_DIR / "eval_wer_grouped.png",
                    Inches(0.6), Inches(2.1), w=Inches(8.5))

    # right table
    rx = Inches(9.3); rw = Inches(3.5)
    add_text(s, "WER · untouched → robust",
             rx, Inches(2.15), rw, Inches(0.4),
             size=12.5, bold=True, color=THEME["secondary"])
    rows = [
        ("Seamless-L",    "18.45%", "14.34%", "22.3% ↓", THEME["primary"]),
        ("Whisper-Large", "28.29%", "19.97%", "29.4% ↓", THEME["accent"]),
        ("Whisper-Med",   "40.44%", "30.64%", "24.2% ↓", THEME["secondary"]),
        ("Wav2Vec2",      "53.52%", "39.67%", "25.9% ↓", THEME["violet"]),
    ]
    y = Inches(2.6)
    for m, a, b, rel, color in rows:
        card = add_card(s, rx, y, rw, Inches(0.95),
                        fill=THEME["card"], line=color, corner=0.08)
        card.line.width = Pt(1.5)
        add_text(s, m, rx + Inches(0.15), y + Inches(0.1),
                 rw - Inches(0.2), Inches(0.32),
                 size=12, bold=True, color=THEME["text"])
        add_text(s, f"{a} → {b}",
                 rx + Inches(0.15), y + Inches(0.42), rw - Inches(0.2),
                 Inches(0.3), size=11, color=THEME["text_dim"])
        add_text(s, rel, rx + Inches(0.15), y + Inches(0.65),
                 rw - Inches(0.2), Inches(0.3), size=12, bold=True,
                 color=color)
        y = y + Inches(1.02)

    add_footer(s, 18, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=120, base_ms=300)


# ============== Slide 19: CER Results ========================================
def slide_results_cer():
    s, bg = page(19)
    add_decoration(s)
    add_header(s, "Results · CER on the same benchmark",
               "Character-level errors fall too — the gains aren't an artefact")

    img = add_image(s, FIG_DIR / "eval_cer_grouped.png",
                    Inches(0.7), Inches(2.2), w=Inches(8.2))

    rx = Inches(9.2); rw = Inches(3.6)
    add_text(s, "Why CER matters",
             rx, Inches(2.2), rw, Inches(0.4),
             size=15, bold=True, color=THEME["secondary"])
    add_text(s,
        "Character Error Rate is robust to tokenisation choices.\n\n"
        "If CER falls in lockstep with WER, we know we're correcting\n"
        "real recognition errors — not gaming a token-boundary metric.",
        rx, Inches(2.65), rw, Inches(2.6),
        size=12, color=THEME["text"], line_spacing=1.45)

    add_text(s, "Highlight",
             rx, Inches(5.3), rw, Inches(0.4),
             size=14, bold=True, color=THEME["primary"])
    add_text(s,
        "Wav2Vec2 CER falls from 36.2% → 30.4%\n"
        "Whisper-Large CER falls from 16.4% → 13.8%",
        rx, Inches(5.7), rw, Inches(1.4),
        size=12, color=THEME["text"], line_spacing=1.45)

    add_footer(s, 19, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=130, base_ms=300)


# ============== Slide 20: Relative Improvement ===============================
def slide_results_relative():
    s, bg = page(20)
    add_decoration(s)
    add_header(s, "Relative WER Reduction",
               "How much of each model's WER did CORAL erase?")

    img = add_image(s, FIG_DIR / "eval_relative_wer_drop.png",
                    Inches(0.7), Inches(2.2), w=Inches(8.5))

    rx = Inches(9.4); rw = Inches(3.4)
    msgs = [
        ("Whisper-Large", "29.4%", THEME["accent"]),
        ("Wav2Vec2-Urdu", "25.9%", THEME["violet"]),
        ("Whisper-Med",   "24.2%", THEME["secondary"]),
        ("Seamless-L",    "22.3%", THEME["primary"]),
    ]
    y = Inches(2.4)
    for m, v, color in msgs:
        add_card(s, rx, y, rw, Inches(0.9), fill=THEME["card"],
                 line=color, corner=0.08)
        add_text(s, v, rx + Inches(0.2), y + Inches(0.1), Inches(1.4),
                 Inches(0.7), size=22, bold=True, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, m, rx + Inches(1.5), y + Inches(0.1),
                 rw - Inches(1.6), Inches(0.7), size=12,
                 color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE, bold=True)
        y = y + Inches(1.0)

    add_text(s,
        "The weaker the underlying ASR, the more headroom CORAL captures —\n"
        "but even Seamless-Large (state-of-the-art) gains 4 absolute WER points.",
        rx, Inches(6.4), rw, Inches(1.0),
        size=11, italic=True, color=THEME["text_dim"], line_spacing=1.4)

    add_footer(s, 20, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 21: Conversational Ablation ============================
def slide_results_ablation():
    s, bg = page(21)
    add_decoration(s)
    add_header(s, "Ablation · Conversational Urdu (500 clips)",
               "C0 → C7  ·  every stage adds independent value")

    img = add_image(s, FIG_DIR / "fyp2_ablation_clean.png",
                    Inches(0.7), Inches(2.1), w=Inches(8.2))
    if img is None:
        img = add_image(s, FIG_DIR / "fyp2_wer_ablation.png",
                        Inches(0.7), Inches(2.1), w=Inches(8.2))

    # right column: per-component gains
    rx = Inches(9.2); rw = Inches(3.6)
    gains = [
        ("Stage 0 normalise", "-1.9", THEME["primary"]),
        ("Stage 1 alignment", "-1.1", THEME["secondary"]),
        ("Stage 2 OOV+BK-tree","-1.2", THEME["accent"]),
        ("Stage 3 voting",    "-1.5", THEME["green"]),
        ("Stage 4 LLM",       "-2.5", THEME["violet"]),
    ]
    add_text(s, "Per-stage WER contribution",
             rx, Inches(2.2), rw, Inches(0.4),
             size=13, bold=True, color=THEME["secondary"])
    y = Inches(2.7)
    for k, v, color in gains:
        add_card(s, rx, y, rw, Inches(0.65), fill=THEME["card"],
                 line=color, corner=0.10)
        add_text(s, k, rx + Inches(0.2), y + Inches(0.12),
                 rw - Inches(1.1), Inches(0.4), size=12,
                 color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE, bold=True)
        add_text(s, v + " pts", rx + rw - Inches(1.1),
                 y + Inches(0.12), Inches(1.0), Inches(0.4),
                 size=13, bold=True, color=color, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.72)

    add_card(s, rx, Inches(6.3), rw, Inches(0.65),
             fill=THEME["primary"], corner=0.10)
    add_text(s, "Total  ·  19.8% → 10.6% WER",
             rx + Inches(0.2), Inches(6.35), rw - Inches(0.2), Inches(0.55),
             size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 21, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=120, base_ms=300)


# ============== Slide 22: Pipeline progression ===============================
def slide_results_progression():
    s, bg = page(22)
    add_decoration(s)
    add_header(s, "Pipeline Progression",
               "WER drops monotonically as each stage is added")

    img = add_image(s, FIG_DIR / "eval_pipeline_progression.png",
                    Inches(1.5), Inches(2.1), w=Inches(10.3))

    add_text(s,
        "Every additional stage in the CORAL pipeline reduces the WER on the conversational benchmark. "
        "The curve is monotone — there is no point at which adding a stage \"undoes\" upstream gain.",
        Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.8),
        size=12.5, color=THEME["text_dim"], italic=True,
        align=PP_ALIGN.CENTER, line_spacing=1.4)

    add_footer(s, 22, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 23: Split/Merge Distribution ============================
def slide_splitmerge():
    s, bg = page(23)
    add_decoration(s)
    add_header(s, "The Split/Merge Phenomenon",
               "36.5% of inter-model events are word-boundary disagreement")

    img = add_image(s, FIG_DIR / "splitmerge_distribution.png",
                    Inches(0.7), Inches(2.2), w=Inches(7.5))

    rx = Inches(8.5); rw = Inches(4.3)
    add_text(s, "What this chart shows",
             rx, Inches(2.2), rw, Inches(0.4),
             size=14, bold=True, color=THEME["secondary"])
    add_text(s,
        "Across all alignment events on the 500-clip Conversational Urdu sample, "
        "SPLIT and MERGE together exceed any single substitution type.",
        rx, Inches(2.65), rw, Inches(1.4),
        size=11.5, color=THEME["text"], line_spacing=1.45)

    add_text(s, "Why ROVER struggles",
             rx, Inches(4.0), rw, Inches(0.4),
             size=14, bold=True, color=THEME["primary"])
    add_text(s,
        "ROVER's WCN edit-distance treats SPLIT as substitution-then-insertion and "
        "MERGE as deletion. The voting signal degrades exactly where it matters most.",
        rx, Inches(4.4), rw, Inches(1.5),
        size=11.5, color=THEME["text"], line_spacing=1.45)

    add_text(s, "CORAL's answer",
             rx, Inches(5.95), rw, Inches(0.4),
             size=14, bold=True, color=THEME["green"])
    add_text(s,
        "Character-level secondary alignment recovers SAME/SPLIT/MERGE/NOISE labels "
        "explicitly — Stage 3 then treats them with appropriate semantics.",
        rx, Inches(6.35), rw, Inches(1.0),
        size=11.5, color=THEME["text"], line_spacing=1.45)

    add_footer(s, 23, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 24: Residual Errors ====================================
def slide_residuals():
    s, bg = page(24)
    add_decoration(s)
    add_header(s, "Residual Error Analysis",
               "What's left after CORAL — and where to dig next")

    img = add_image(s, FIG_DIR / "residual_errors.png",
                    Inches(0.7), Inches(2.2), w=Inches(7.5))

    rx = Inches(8.5); rw = Inches(4.3)
    items = [
        ("27.7%", "Proper nouns", "Named-entity gazetteer needed", THEME["primary"]),
        ("21.5%", "Code-switching", "Stage 0 strips English — by design", THEME["accent"]),
        ("13.4%", "Dialectal forms", "Limited corpus coverage",      THEME["secondary"]),
        ("9.8%",  "Numerals & dates","Acoustic ambiguity",           THEME["green"]),
        ("7.2%",  "Over-correction", "Top-1 BK-tree too greedy",     THEME["violet"]),
        ("20.4%", "Other / acoustic","Genuine recognition failure",  THEME["text_dim"]),
    ]
    y = Inches(2.2)
    for v, k, why, color in items:
        c = add_card(s, rx, y, rw, Inches(0.7),
                     fill=THEME["card"], line=color, corner=0.08)
        add_text(s, v, rx + Inches(0.2), y + Inches(0.15),
                 Inches(1.0), Inches(0.45), size=16, bold=True,
                 color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, k, rx + Inches(1.2), y + Inches(0.08),
                 rw - Inches(1.3), Inches(0.3),
                 size=11.5, bold=True, color=THEME["text"])
        add_text(s, why, rx + Inches(1.2), y + Inches(0.38),
                 rw - Inches(1.3), Inches(0.35),
                 size=9.5, color=THEME["text_dim"], italic=True)
        y = y + Inches(0.78)

    add_footer(s, 24, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=120, base_ms=300)


# ============== Slide 25: Comparison to SOTA =================================
def slide_comparison():
    s, bg = page(25)
    add_decoration(s)
    add_header(s, "Head-to-Head Comparison",
               "CORAL vs ROVER and the closest concurrent system")

    headers = ["System", "Conv. Urdu WER", "Approach", "Cost"]
    rows = [
        ("Whisper-Large (raw)",    "19.8%", "Single-model baseline",                 "1× ASR",          THEME["text_dim"]),
        ("ROVER (English-style)",  "18.2%", "WCN voting on raw hypotheses",          "1× align",        THEME["accent"]),
        ("Multi-ASR + SpeechLLM",  "11.3%", "Heavy SpeechLLM audio re-inference",    "+ SpeechLLM pass",THEME["secondary"]),
        ("CORAL  (this work)",     "10.6%", "5-stage post-processing · text-only",   "+ 200 ms / utt",  THEME["primary"]),
    ]
    # Header row
    col_w = [Inches(3.5), Inches(2.0), Inches(4.4), Inches(2.6)]
    x_pos = [Inches(0.7), Inches(4.2), Inches(6.2), Inches(10.6)]
    add_card(s, Inches(0.7), Inches(2.2), Inches(12.5), Inches(0.55),
             fill=THEME["card"], line=THEME["secondary"])
    for h, x_, w in zip(headers, x_pos, col_w):
        add_text(s, h, x_ + Inches(0.15), Inches(2.3),
                 w - Inches(0.3), Inches(0.4), size=12, bold=True,
                 color=THEME["secondary"])

    y = Inches(2.9)
    for row in rows:
        name, wer, approach, cost, color = row
        c = add_card(s, Inches(0.7), y, Inches(12.5), Inches(0.75),
                     fill=THEME["card"], line=color, corner=0.05)
        add_text(s, name, x_pos[0] + Inches(0.15), y + Inches(0.1),
                 col_w[0] - Inches(0.3), Inches(0.55), size=13,
                 bold=True, color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, wer, x_pos[1] + Inches(0.15), y + Inches(0.1),
                 col_w[1] - Inches(0.3), Inches(0.55), size=15,
                 bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, approach, x_pos[2] + Inches(0.15), y + Inches(0.1),
                 col_w[2] - Inches(0.3), Inches(0.55), size=11.5,
                 color=THEME["text"], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cost, x_pos[3] + Inches(0.15), y + Inches(0.1),
                 col_w[3] - Inches(0.3), Inches(0.55), size=11,
                 color=THEME["text_dim"], anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.85)

    add_text(s,
        "CORAL beats Multi-ASR + SpeechLLM by 0.7 absolute WER points without running "
        "an additional audio-input LLM — entirely as text-level post-processing.",
        Inches(0.7), Inches(6.6), Inches(12.5), Inches(0.7),
        size=12, italic=True, color=THEME["accent"], align=PP_ALIGN.CENTER)

    add_footer(s, 25, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 26: Key Achievements ===================================
def slide_achievements():
    s, bg = page(26)
    add_decoration(s)
    add_header(s, "Key Achievements", "Seven concrete wins from CORAL")

    items = [
        ("Novel Algorithm",
         "First Urdu-aware split-merge alignment with weighted Levenshtein.",
         THEME["primary"]),
        ("SOTA Beat",
         "10.6% WER on Conversational Urdu — beats Multi-ASR+SpeechLLM (11.3%).",
         THEME["accent"]),
        ("Zero Retraining",
         "All gains delivered at inference time. No acoustic-model fine-tune.",
         THEME["secondary"]),
        ("Reproducible Eval",
         "Open TSV with 16 (model × config) WER+CER rows.",
         THEME["green"]),
        ("Deployed System",
         "FastAPI + Next.js, with Kaggle-GPU registry & ngrok fan-out.",
         THEME["blue"]),
        ("Compositional",
         "Every stage adds independent, additive WER reduction.",
         THEME["violet"]),
    ]
    x = Inches(0.7); y = Inches(2.3)
    w = Inches(6.0); h = Inches(1.45); gap_y = Inches(0.15)
    for i, (head, body, color) in enumerate(items):
        cx = x if i % 2 == 0 else x + w + Inches(0.25)
        cy = y + (i // 2) * (h + gap_y)
        c = add_card(s, cx, cy, w, h, fill=THEME["card"],
                     line=color, corner=0.07)
        c.line.width = Pt(2.0)
        # left coloured bar
        bar = add_rect(s, cx, cy, Inches(0.18), h, color,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.adjustments[0] = 0.4
        add_text(s, head, cx + Inches(0.35), cy + Inches(0.15),
                 w - Inches(0.5), Inches(0.45),
                 size=15, bold=True, color=THEME["text"])
        add_text(s, body, cx + Inches(0.35), cy + Inches(0.65),
                 w - Inches(0.5), Inches(0.75),
                 size=11.5, color=THEME["text_dim"], line_spacing=1.4)

    add_footer(s, 26, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=130, base_ms=300)


# ============== Slide 27: Future Work ========================================
def slide_future():
    s, bg = page(27)
    add_decoration(s)
    add_header(s, "Future Work", "10 directions for follow-up research")

    items = [
        ("1", "Confidence-weighted voting",      THEME["primary"]),
        ("2", "Code-switching-aware Stage 0",    THEME["accent"]),
        ("3", "Named-entity gazetteer",          THEME["secondary"]),
        ("4", "LLM top-K OOV re-rank",           THEME["violet"]),
        ("5", "Split/Merge in voting logic",     THEME["green"]),
        ("6", "Learned cross-arch confidence",   THEME["orange"]),
        ("7", "Larger conversational eval",      THEME["blue"]),
        ("8", "Real-time streaming variant",     THEME["primary"]),
        ("9", "Other low-resource SA langs",     THEME["accent"]),
        ("10", "Open-source release",            THEME["secondary"]),
    ]
    x = Inches(0.6); w = Inches(2.45); h = Inches(1.95)
    gap = Inches(0.05)
    cols, rows = 5, 2
    for i, (n, lbl, color) in enumerate(items):
        cx = x + (i % cols) * (w + gap)
        cy = Inches(2.3) + (i // cols) * (h + Inches(0.2))
        c = add_card(s, cx, cy, w, h, fill=THEME["card"],
                     line=color, corner=0.07)
        c.line.width = Pt(2.0)
        add_text(s, n, cx, cy + Inches(0.25), w, Inches(0.6),
                 size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, lbl, cx + Inches(0.15), cy + Inches(0.95),
                 w - Inches(0.3), Inches(0.95),
                 size=11.5, color=THEME["text"], align=PP_ALIGN.CENTER,
                 line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 27, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=120, base_ms=280)


# ============== Slide 28: Section divider — Demo =============================
def slide_demo_divider():
    s, bg = page(28)
    # giant decorative shapes
    b1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(-2), Inches(-2), Inches(6), Inches(6))
    b1.line.fill.background(); b1.fill.solid()
    b1.fill.fore_color.rgb = THEME["primary"]; b1.shadow.inherit=False
    _set_transparency(b1, 60)
    b2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(9), Inches(3), Inches(6), Inches(6))
    b2.line.fill.background(); b2.fill.solid()
    b2.fill.fore_color.rgb = THEME["secondary"]; b2.shadow.inherit=False
    _set_transparency(b2, 70)

    add_text(s, "LIVE DEMO", Inches(0.7), Inches(2.5), Inches(12),
             Inches(0.6), size=18, bold=True, color=THEME["secondary"])
    add_text(s, "From audio to corrected Urdu.\nIn 4 user-facing passes.",
             Inches(0.7), Inches(3.0), Inches(12), Inches(2.5),
             size=60, bold=True, color=THEME["text"], line_spacing=1.1)

    add_text(s, "1.  Speech upload   →   2. Alignment scan   →   3. OOV sieve animation   →   4. Voting / LLM diff",
             Inches(0.7), Inches(5.7), Inches(12), Inches(0.5),
             size=14, color=THEME["text_dim"])

    set_slide_transition(s, "push")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id),
                                                        int(b1.shape_id),
                                                        int(b2.shape_id)}),
                           delay_ms=180, base_ms=400)


# ============== Slide 29: Frontend workflow ==================================
def slide_frontend():
    s, bg = page(29)
    add_decoration(s)
    add_header(s, "Frontend — R&D Product",
               "Five-page Next.js platform · dark cinematic theme")

    # ── Top row: site map ──────────────────────────────────────────────────
    pages = [
        ("/",         "Landing",   "Hero · waveform\nKPIs · pipeline\nresults · stack", THEME["primary"]),
        ("/pipeline", "Pipeline",  "Stage 00 → 04\nexamples + metrics\nflow diagram",   THEME["secondary"]),
        ("/research", "Research",  "Hypothesis\nablation chart\nresiduals · future",   THEME["accent"]),
        ("/team",     "Team",      "Members + supervisors\nFAST School of Computing\nphotos + socials", THEME["violet"]),
        ("/app",      "Demo",      "3-pass interactive\nfile · speech · manual\nLLM refinement",     THEME["green"]),
    ]
    x = Inches(0.5); w = Inches(2.46); gap = Inches(0.12)
    for i, (route, head, body, color) in enumerate(pages):
        c = add_card(s, x, Inches(2.2), w, Inches(2.45),
                     fill=THEME["card"], line=color, corner=0.07)
        c.line.width = Pt(2.0)
        add_pill(s, route, x + Inches(0.15), Inches(2.30),
                 w - Inches(0.3), Inches(0.35),
                 color, RGBColor(0xFF, 0xFF, 0xFF), size=10)
        add_text(s, head, x, Inches(2.75), w, Inches(0.4),
                 size=15, bold=True, color=THEME["text"], align=PP_ALIGN.CENTER)
        add_text(s, body, x + Inches(0.2), Inches(3.20),
                 w - Inches(0.4), Inches(1.30),
                 size=10.5, color=THEME["text_dim"],
                 align=PP_ALIGN.CENTER, line_spacing=1.4)
        if i < len(pages) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + w - Inches(0.02),
                                      Inches(3.30), Inches(0.16),
                                      Inches(0.30))
            arr.fill.solid(); arr.fill.fore_color.rgb = THEME["text_dim"]
            arr.line.fill.background(); arr.shadow.inherit=False
        x = x + w + gap

    # ── Bottom row: interactive demo passes ────────────────────────────────
    add_text(s, "Interactive demo — 3-pass workflow",
             Inches(0.5), Inches(4.85), Inches(12), Inches(0.35),
             size=11, bold=True, color=THEME["secondary"])
    passes = [
        ("Pass 1", "Alignment",  "Input audio (mic / MP3) or\npre-aligned TSV → /align", THEME["primary"]),
        ("Pass 2", "OOV Sieve",  "Animated 4-phase BK-tree\nscan + n-gram re-rank",      THEME["accent"]),
        ("Pass 3", "Correction", "Voting · LLM (GPT-OSS)\nside-by-side diff",             THEME["green"]),
    ]
    x = Inches(0.5); w = Inches(4.16); gap = Inches(0.15)
    for i, (n, head, body, color) in enumerate(passes):
        c = add_card(s, x, Inches(5.25), w, Inches(1.50),
                     fill=THEME["card"], line=color, corner=0.06)
        add_pill(s, n, x + Inches(0.2), Inches(5.35),
                 Inches(0.8), Inches(0.3),
                 color, RGBColor(0xFF, 0xFF, 0xFF), size=9)
        add_text(s, head, x + Inches(1.1), Inches(5.32),
                 Inches(2.8), Inches(0.35),
                 size=13, bold=True, color=THEME["text"])
        add_text(s, body, x + Inches(0.2), Inches(5.70),
                 w - Inches(0.4), Inches(1.05),
                 size=10.5, color=THEME["text_dim"], line_spacing=1.35)
        x = x + w + gap

    add_text(s,
        "Dark cinematic theme  •  animated waveform hero  •  aurora gradients  •  particle field  •  staggered reveals  •  localStorage state",
        Inches(0.5), Inches(6.90), Inches(12.5), Inches(0.3),
        size=10, italic=True, color=THEME["text_dim"], align=PP_ALIGN.CENTER)

    add_footer(s, 29, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=140, base_ms=300)


# ============== Slide 30: Deliverables =======================================
def slide_deliverables():
    s, bg = page(30)
    add_decoration(s)
    add_header(s, "Deliverables", "Everything shipped with CORAL")

    items = [
        ("Pipeline",     "5-stage Python pipeline\nFastAPI + Docker",        THEME["primary"]),
        ("Frontend",     "Next.js + React app\nFile + Speech modes",        THEME["accent"]),
        ("Kaggle ASR",   "GPU notebooks (Whisper-L, Seamless-L)\nself-register via ngrok", THEME["secondary"]),
        ("Registry",     "60-s heartbeats\n300-s eviction",                  THEME["green"]),
        ("Evaluation",   "WER + CER · 4 configs × 4 models\nensemble_correction_eval.tsv", THEME["violet"]),
        ("Data",         "BK-tree.joblib + n-gram Parquet\nHF-hosted",       THEME["blue"]),
        ("Docs",         "97-page LaTeX report\n+ FYP-1 & FYP-2 PDFs",       THEME["orange"]),
        ("Repository",   "Open code + permissive licence",                   THEME["primary"]),
    ]
    x = Inches(0.6); y = Inches(2.3)
    w = Inches(3.05); h = Inches(2.0); gap = Inches(0.05)
    for i, (k, v, color) in enumerate(items):
        cx = x + (i % 4) * (w + gap)
        cy = y + (i // 4) * (h + Inches(0.15))
        c = add_card(s, cx, cy, w, h,
                     fill=THEME["card"], line=color, corner=0.08)
        c.line.width = Pt(2.0)
        add_pill(s, k, cx + Inches(0.2), cy + Inches(0.2),
                 w - Inches(0.4), Inches(0.4),
                 color, RGBColor(0xFF, 0xFF, 0xFF), size=12)
        add_text(s, v, cx + Inches(0.2), cy + Inches(0.75),
                 w - Inches(0.4), h - Inches(0.9),
                 size=11.5, color=THEME["text"], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

    add_footer(s, 30, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=120, base_ms=300)


# ============== Slide 31: Conclusion =========================================
def slide_conclusion():
    s, bg = page(31)
    add_decoration(s)
    add_header(s, "Conclusion", "What CORAL proves")

    statements = [
        ("Word-boundary disagreement is the dominant inter-model error mode in Urdu ASR.",
         THEME["primary"]),
        ("It can be solved algorithmically — Stages 0-3 cut WER by 6.2 absolute points alone.",
         THEME["secondary"]),
        ("An LLM is most useful as the LAST refinement step, not the primary fusion engine.",
         THEME["violet"]),
        ("Compositional design wins: every stage's gain is additive — none cancel.",
         THEME["green"]),
        ("Hypothesis confirmed: 5 levers + no fine-tuning  →  10.6% Conv. Urdu WER.",
         THEME["accent"]),
    ]
    y = Inches(2.3)
    for txt, color in statements:
        c = add_card(s, Inches(0.7), y, Inches(12.0), Inches(0.85),
                     fill=THEME["card"], line=color, corner=0.10)
        # left bullet
        dot = add_rect(s, Inches(0.85), y + Inches(0.27), Inches(0.32),
                       Inches(0.32), color, shape=MSO_SHAPE.OVAL)
        add_text(s, txt, Inches(1.3), y + Inches(0.15), Inches(11.2),
                 Inches(0.55), size=14, color=THEME["text"],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        y = y + Inches(0.95)

    add_footer(s, 31, TOTAL_SLIDES)
    set_slide_transition(s, "fade")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id)}),
                           delay_ms=180, base_ms=350)


# ============== Slide 32: Thank You / Q&A ===================================
def slide_thank_you():
    s, bg = page(32)
    # huge decoration
    b1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(-3), Inches(2), Inches(8), Inches(8))
    b1.line.fill.background(); b1.fill.solid()
    b1.fill.fore_color.rgb = THEME["primary"]; b1.shadow.inherit=False
    _set_transparency(b1, 65)
    b2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(8), Inches(-2), Inches(8), Inches(8))
    b2.line.fill.background(); b2.fill.solid()
    b2.fill.fore_color.rgb = THEME["secondary"]; b2.shadow.inherit=False
    _set_transparency(b2, 70)

    add_text(s, "Thank you.",
             Inches(0.7), Inches(2.4), Inches(12), Inches(1.4),
             size=88, bold=True, color=THEME["text"])
    add_text(s, "Questions, comments, or constructive disagreement — all welcome.",
             Inches(0.7), Inches(4.0), Inches(12), Inches(0.6),
             size=18, color=THEME["text"], italic=True)

    # contact band
    c = add_card(s, Inches(0.7), Inches(5.4), Inches(12.0), Inches(1.4),
                 fill=THEME["card"], line=THEME["primary"])
    add_text(s, "CORAL  •  CORAL-Urdu-ASR repository",
             Inches(1.0), Inches(5.6), Inches(11.5), Inches(0.45),
             size=14, bold=True, color=THEME["primary"])
    add_text(s,
        "Ali Irfan  21I-2572   •   Nouman Hafeez  21I-0416   •   Rafay Khattak  21I-0423\n"
        "Dept. of Computer Science · FAST School of Computing · FAST-NUCES Islamabad\n"
        "Supervised by Ms. Kainat Iqbal · co-supervised by Ms. Saira Qamar  (AI & Data Science)",
        Inches(1.0), Inches(5.95), Inches(11.5), Inches(1.05),
        size=12, color=THEME["text"], line_spacing=1.4)

    set_slide_transition(s, "push")
    inject_fade_animations(s, collect_anim(s, exclude={int(bg.shape_id),
                                                        int(b1.shape_id),
                                                        int(b2.shape_id)}),
                           delay_ms=200, base_ms=450)


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------
slide_title()
slide_team()
slide_agenda()
slide_problem()
slide_hypothesis()
slide_evolution()
slide_dropped()
slide_pipeline()
slide_stage0()
slide_stage1()
slide_stage2()
slide_stage3()
slide_stage4()
slide_architecture()
slide_deployment()
slide_techstack()
slide_eval_setup()
slide_results_readspeech()
slide_results_cer()
slide_results_relative()
slide_results_ablation()
slide_results_progression()
slide_splitmerge()
slide_residuals()
slide_comparison()
slide_achievements()
slide_future()
slide_demo_divider()
slide_frontend()
slide_deliverables()
slide_conclusion()
slide_thank_you()

prs.save(OUT)
print(f"OK saved {OUT}  ({OUT.stat().st_size/1024:.1f} KB)")
